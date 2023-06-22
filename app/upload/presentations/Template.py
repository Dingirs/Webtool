import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
prs = Presentation()
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Hot Chocolate Bomb In Full Color Gift Box"
content = slide.placeholders[1]
content.text = "Item number: CUKID-OBYBW\n" \
"\n" \
"Get in on the hottest item of the Holiday Season! Our BRAND NEW Hot Chocolate Bomb is not only delicious but an amazing food experience! Filled with delicious marshmallows, simply place the chocolate bomb into your mug, pour 10 oz of hot milk over it and watch the bomb unravel before your eyes! Fast and easy, this hot cocoa recipe is sure to make a cold winter\'s day much better. Check out our video: https://youtu.be/aF80JkdF2Xo. 3\" L x 3\" W x 3\" H\n" \
"\n" \
"Colors: White\n" \
"\n" \
"Decoration Information: 1\" W x 2\" H; Box\n" \
"\n" \
"Price Table:\n" \
"| Qty  | 50    | 100  | 250  | 500  |\n" \
"|------|-------|------|------|------|\n" \
"| Price| $6.55 | $6.24 | $6.06 | $5.88 |\n" \
"\n" \
"Price Includes: 1 color;1 location\n" \
"\n" \
"Packaging and Delivery: Bulk. 24 units per carton. 10 lbs. per carton. Normal production time is 10 working days."
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "15 Oz. Campfire Mug With Mug Stuffer"
content = slide.placeholders[1]
content.text = "Item number: QTLEK-NMBOT\n" \
"\n" \
"Includes 1-Color/1-Location Imprint On #7133 - 15 Oz. Campfire Mug And Stuffer Bag With Your Choice Of Filler. Stuffer Package Will Not Be Imprinted. Retro Granite Design. Meets FDA Requirements. Hand Wash Recommended. Not Recommended for Commercial Use. 3 1/2\" H\n" \
"\n" \
"Colors: Ms22 Fill In #7133 Campfire\n" \
"\n" \
"Price Table:\n" \
"| Qty  | 144  | 288  | 576  | 1008 |\n" \
"|------|------|------|------|------|\n" \
"| Price| $9.81 | $8.53 | $7.42 | $6.45 |\n" \
"\n" \
"Price Includes: 1 Color;1 Location\n" \
"\n" \
"Packaging and Delivery: Bulk. 36 units per carton. 44 lbs. per carton. Normal production time is 3 to 10 working days."

prs.save("/code/app/upload/presentations/Template2.pptx")